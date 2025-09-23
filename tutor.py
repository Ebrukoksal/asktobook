# tutor.py

import os
import re
from typing import List, Literal, Optional

import chromadb
from openai import OpenAI
from sentence_transformers import SentenceTransformer

try:
    from pptx import Presentation
except Exception:
    Presentation = None  # optional dependency guard


def detect_page_or_slide_request(query: str) -> Optional[tuple[str, int]]:
    """
    Detect if query explicitly asks about a specific page/slide.
    Returns ("page"/"slide", number) or None.
    Turkish and English variants supported.
    """
    match = re.search(r"(?:\b(?:sayfa|page)\b)\s*(\d+)", query, re.IGNORECASE)
    if match:
        return "page", int(match.group(1))
    match = re.search(r"(?:\b(?:slayt|slide)\b)\s*(\d+)", query, re.IGNORECASE)
    if match:
        return "slide", int(match.group(1))
    return None


class TextbookTutor:
    def __init__(self, model_name: str = "gpt-4.1-nano"):
        self.model_name = model_name
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        # HuggingFace embedding
        self.embedder = SentenceTransformer("all-MiniLM-L6-v2")

        # ChromaDB setup
        self.chroma_client = chromadb.Client()
        try:
            self.collection = self.chroma_client.get_collection(name="textbook")
        except Exception:
            self.collection = self.chroma_client.create_collection(name="textbook")

        self.sections: List[str] = []
        self.label_prefix: Literal["Page", "Slide"] = "Page"

    def semantic_sections(self, text: str) -> List[str]:
        """Generate semantic sections from text."""
        return [s.strip() for s in text.split("\n\n") if s.strip()]

    def extract_pptx_slides(self, file_path: str) -> List[str]:
        """Extract text per slide from a PPTX file."""
        if Presentation is None:
            raise RuntimeError("python-pptx not installed. Run `pip install python-pptx`.")
        prs = Presentation(file_path)
        slides_text: List[str] = []
        for slide in prs.slides:
            chunks = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    paragraphs = []
                    for p in shape.text_frame.paragraphs:
                        runs = [r.text for r in getattr(p, "runs", [])] or [p.text]
                        paragraphs.append("".join(runs))
                    if paragraphs:
                        chunks.append("\n".join(paragraphs))
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_paras = []
                for p in slide.notes_slide.notes_text_frame.paragraphs:
                    runs = [r.text for r in getattr(p, "runs", [])] or [p.text]
                    notes_paras.append("".join(runs))
                if notes_paras:
                    chunks.append("\n[Notes]\n" + "\n".join([p for p in notes_paras if p]))
            slides_text.append("\n".join([c for c in chunks if c]).strip())
        return [s if s else "" for s in slides_text]

    def load_textbook(self, pages: List[str], label_prefix: Literal["Page", "Slide"] = "Page"):
        """
        Load a list of page/slide texts (already extracted).
        label_prefix: "Page" (PDF) or "Slide" (PPTX).
        """
        self.label_prefix = label_prefix
        full_text = "\n".join(pages)
        sections = self.semantic_sections(full_text)
        self.sections = sections

        # Recreate collection to avoid id collisions on reload
        try:
            self.chroma_client.delete_collection("textbook")
        except Exception:
            pass
        self.collection = self.chroma_client.create_collection(name="textbook")

        for i, section in enumerate(sections):
            emb = self.embedder.encode([section])[0].tolist()
            self.collection.add(
                ids=[str(i)],
                embeddings=[emb],
                documents=[section],
                metadatas=[{"page": i + 1, "label_prefix": self.label_prefix}],
            )
        print(f"✅ Loaded {len(sections)} semantic sections into ChromaDB")

    def _find_relevant_chunks(self, query: str, top_k: int = 2):
        """Retrieve top-k relevant chunks from ChromaDB."""
        query_emb = self.embedder.encode([query])[0].tolist()
        results = self.collection.query(query_embeddings=[query_emb], n_results=top_k)
        return results

    def answer(self, query: str, top_k: int = 2):
        """
        Generate answer using either:
        - Direct page/slide addressing (if query references one), or
        - Retrieved top-k chunks from ChromaDB.
        """
        # 1) Direct page/slide addressing
        page_or_slide = detect_page_or_slide_request(query)
        if page_or_slide and self.sections:
            kind, num = page_or_slide
            if 1 <= num <= len(self.sections):
                context_text = f"{self.label_prefix} {num}:\n{self.sections[num - 1]}"
                system_prompt = f"""
You are an educational assistant.
Summarize or explain ONLY the content of {self.label_prefix.lower()} {num}.
Do NOT add external knowledge.

Content:
{context_text}

Instructions:
1. Answer based only on the given {self.label_prefix.lower()}.
2. Mention {self.label_prefix} {num} explicitly.
3. Keep the answer concise and accurate.
"""
                try:
                    response = self.client.chat.completions.create(
                        model=self.model_name,
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": query},
                        ],
                        temperature=0.2,
                        max_tokens=500,
                    )
                    answer = response.choices[0].message.content.strip()
                    return {"answer": answer, "pages": [num], "label_prefix": self.label_prefix}
                except Exception as e:
                    return {"answer": f"⚠️ Error: {str(e)}", "pages": [], "label_prefix": self.label_prefix}
            else:
                return {
                    "answer": f"⚠️ {self.label_prefix} {num} not found in this document.",
                    "pages": [],
                    "label_prefix": self.label_prefix,
                }

        # 2) Retrieval-augmented answering (fallback)
        results = self._find_relevant_chunks(query, top_k=top_k)
        if not results.get("documents"):
            return {"answer": "⚠️ No relevant information found.", "pages": [], "label_prefix": self.label_prefix}

        docs = results["documents"][0]
        pages = [meta["page"] for meta in results["metadatas"][0]]
        prefixes = [meta.get("label_prefix", self.label_prefix) for meta in results["metadatas"][0]]

        context_chunks = [f"{pfx} {pnum}:\n{doc}" for pfx, pnum, doc in zip(prefixes, pages, docs)]
        context_text = "\n\n".join(context_chunks)

        system_prompt = f"""
You are an educational assistant.
Use ONLY the following textbook content to answer the question.
Do NOT add external knowledge.

Textbook content:
{context_text}

Instructions:
1. Answer based only on the given pages/slides.
2. Mention which page(s)/slide(s) the answer came from.
3. Keep the answer concise and accurate.
"""
        try:
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": query},
                ],
                temperature=0.2,
                max_tokens=500,
            )
            answer = response.choices[0].message.content.strip()
            return {"answer": answer, "pages": pages, "label_prefix": self.label_prefix}
        except Exception as e:
            return {"answer": f"⚠️ Error: {str(e)}", "pages": [], "label_prefix": self.label_prefix}
