# streamlit.py

import os
import re
from typing import List

import numpy as np
import PyPDF2
import streamlit as st
from openai import OpenAI
from sklearn.feature_extraction.text import TfidfVectorizer

try:
    from pptx import Presentation
except Exception:
    Presentation = None

st.title("Ask Questions About Your Book (PDF & PPTX)")


# -------- Helpers --------
def detect_page_or_slide_request(query: str):
    """
    Detect if query explicitly asks about a specific page/slide.
    Returns ("page"/"slide", number) or None.
    Turkish and English variants supported.
    """
    match = re.search(r"(?:\b(?:sayfa|page)\b)\s*(\d+)", query, re.IGNORECASE)
    if match:
        return "page", int(match.group(1))
    match = re.search(r"(?:\b(?:slayt|slide)\b)\s*(\d+)", query, re.IGNORECASE)
    if match:
        return "slide", int(match.group(1))
    return None


def extract_pdf_pages(file) -> List[str]:
    reader = PyPDF2.PdfReader(file)
    return [page.extract_text() or "" for page in reader.pages]


def extract_pptx_slides(file) -> List[str]:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. Run `pip install python-pptx`.")
    prs = Presentation(file)

    def _collect_shape_text(shape):
        chunks = []

        # 1) Plain text frames (text boxes, placeholders, titles)
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            paragraphs = []
            for p in shape.text_frame.paragraphs:
                runs = [r.text for r in getattr(p, "runs", [])] or [p.text]
                paragraphs.append("".join(runs))
            if paragraphs:
                chunks.append("\n".join(paragraphs))

        # 2) Tables (optional but common in slides)
        if hasattr(shape, "table") and shape.table:
            rows_text = []
            for row in shape.table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text]
                if row_text:
                    rows_text.append("\t".join(row_text))
            if rows_text:
                chunks.append("\n".join(rows_text))

        # 3) Grouped shapes (recurse)
        if hasattr(shape, "shapes"):  # grouped shape container
            for shp in shape.shapes:
                chunks.extend(_collect_shape_text(shp))

        return chunks

    slides_text: List[str] = []
    for slide in prs.slides:
        chunks = []
        for shape in slide.shapes:
            chunks.extend(_collect_shape_text(shape))

        # Notes (if any)
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_paras = []
            for p in slide.notes_slide.notes_text_frame.paragraphs:
                runs = [r.text for r in getattr(p, "runs", [])] or [p.text]
                notes_paras.append("".join(runs))
            if notes_paras:
                chunks.append("\n[Notes]\n" + "\n".join([p for p in notes_paras if p]))

        slides_text.append("\n".join([c for c in chunks if c]).strip())

    return [s if s else "" for s in slides_text]



def build_index(pages):
    vectorizer = TfidfVectorizer(stop_words="english")
    embeddings = vectorizer.fit_transform(pages)
    return vectorizer, embeddings


def find_relevant_pages(query, vectorizer, embeddings, top_k=2):
    query_vec = vectorizer.transform([query])
    scores = np.dot(query_vec, embeddings.T).toarray()[0]
    return np.argsort(scores)[::-1][:top_k]


# -------- API Key --------
user_api_key = st.text_input("Enter your OpenAI API key:", type="password")
if not user_api_key:
    st.warning("Please enter your OpenAI API key to continue.")
    st.stop()


# -------- Book selection / upload --------
book_options = [
    "Book 1: Interprocess Communications in LINUX",
    "Book 2: C# Network Programming",
    "Upload your own (PDF/PPTX)",
]
book_choice = st.radio("Choose a book or upload your own:", book_options)

label_prefix = "Page"  # default for PDF
if book_choice == "Upload your own (PDF/PPTX)":
    uploaded_file = st.file_uploader("Upload a PDF or PPTX", type=["pdf", "pptx"])
    if not uploaded_file:
        st.info("Please upload a file to get started.")
        st.stop()
    filename = uploaded_file.name.lower()
    if filename.endswith(".pdf"):
        pdf_text_pages = extract_pdf_pages(uploaded_file)
        label_prefix = "Page"
    elif filename.endswith(".pptx"):
        pdf_text_pages = extract_pptx_slides(uploaded_file)
        label_prefix = "Slide"
    else:
        st.error("Unsupported file type.")
        st.stop()
else:
    if book_choice == "Book 1: Interprocess Communications in LINUX":
        with open("books/computer network programming kitap 1.pdf", "rb") as f:
            pdf_text_pages = extract_pdf_pages(f)
        label_prefix = "Page"
    elif book_choice == "Book 2: C# Network Programming":
        with open("books/computer network programming kitap 2.pdf", "rb") as f:
            pdf_text_pages = extract_pdf_pages(f)
        label_prefix = "Page"

vectorizer, embeddings = build_index(pdf_text_pages)

# -------- Chat state --------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
st.session_state["label_prefix"] = label_prefix  # keep fresh


# -------- Chat input --------
prompt = st.chat_input("Ask about your book")
if prompt:
    st.session_state.chat_history.append({"question": prompt, "answer": None, "pages_used": []})


# -------- Chat UI --------
chat_area = st.container()

with chat_area:
    if st.session_state.chat_history:
        st.subheader("Chat History")
        for i, entry in enumerate(st.session_state.chat_history):
            with st.chat_message("user"):
                st.markdown(f"**You:** {entry['question']}")
            if entry["answer"]:
                with st.chat_message("assistant"):
                    st.markdown(f"**Answer:** {entry['answer']}")
                    if entry["pages_used"]:
                        st.markdown("**Sources:**")
                        for idx in entry["pages_used"]:
                            if 0 <= idx - 1 < len(pdf_text_pages):
                                section_text = pdf_text_pages[idx - 1].strip()
                                snippet = section_text[:300] + ("..." if len(section_text) > 300 else "")
                                with st.expander(f"{st.session_state['label_prefix']} {idx}"):
                                    st.write(snippet)


# -------- Answer generation --------
for i, entry in enumerate(st.session_state.chat_history):
    if entry["answer"] is None:
        with st.spinner("Thinking..."):
            client = OpenAI(api_key=user_api_key)

            # --- Choose retrieval path: direct page/slide vs semantic ---
            page_or_slide = detect_page_or_slide_request(entry["question"])
            if page_or_slide:
                # Directly use the requested page/slide
                kind, num = page_or_slide
                if 1 <= num <= len(pdf_text_pages):
                    context_text = f"{st.session_state['label_prefix']} {num}:\n{pdf_text_pages[num-1]}"
                    messages = [
                        {
                            "role": "system",
                            "content": (
                                f"You are a helpful assistant. Summarize or explain ONLY the content of "
                                f"{st.session_state['label_prefix'].lower()} {num} from the textbook."
                            ),
                        },
                        {"role": "user", "content": f"{context_text}\n\nQuestion: {entry['question']}\nAnswer:"},
                    ]
                    pages_used = [num]
                else:
                    st.session_state.chat_history[i]["answer"] = (
                        f"⚠️ {st.session_state['label_prefix']} {num} not found in this document."
                    )
                    st.session_state.chat_history[i]["pages_used"] = []
                    with chat_area:
                        with st.chat_message("assistant"):
                            st.markdown(f"**Answer:** {st.session_state.chat_history[i]['answer']}")
                    break
            else:
                # Fallback: semantic retrieval
                relevant_pages_idx = find_relevant_pages(entry["question"], vectorizer, embeddings, top_k=2)
                context_parts = [
                    f"{st.session_state['label_prefix']} {j+1}:\n{pdf_text_pages[j]}" for j in relevant_pages_idx
                ]
                context_text = "\n\n".join(context_parts)

                messages = [
                    {
                        "role": "system",
                        "content": (
                            "You are a helpful assistant. Use ONLY the information from the provided textbook "
                            "pages/slides to answer. Do NOT add external knowledge. Mention which page(s)/slide(s) "
                            "the answer came from."
                        ),
                    }
                ]
                for prev in st.session_state.chat_history[:i]:
                    messages.append({"role": "user", "content": prev["question"]})
                    messages.append({"role": "assistant", "content": prev["answer"]})

                messages.append(
                    {"role": "user", "content": f"Textbook content:\n{context_text}\n\nQuestion: {entry['question']}\nAnswer:"}
                )
                pages_used = [j + 1 for j in relevant_pages_idx]

            # --- Call the model ---
            response = client.chat.completions.create(
                model="gpt-4.1-nano",
                messages=messages,
                temperature=0.2,
                max_tokens=500,
            )

            answer_text = response.choices[0].message.content.strip()

            # Update state
            st.session_state.chat_history[i]["answer"] = answer_text
            st.session_state.chat_history[i]["pages_used"] = pages_used

            # Partial render
            with chat_area:
                with st.chat_message("assistant"):
                    st.markdown(f"**Answer:** {answer_text}")
                    if pages_used:
                        st.markdown("**Sources:**")
                        for idx in pages_used:
                            if 0 <= idx - 1 < len(pdf_text_pages):
                                section_text = pdf_text_pages[idx - 1].strip()
                                snippet = section_text[:300] + ("..." if len(section_text) > 300 else "")
                                with st.expander(f"{st.session_state['label_prefix']} {idx}"):
                                    st.write(snippet)
        break
